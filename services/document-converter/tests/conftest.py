from __future__ import annotations

from io import BytesIO

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation


@pytest.fixture
def docx_bytes() -> bytes:
    document = Document()
    document.add_heading("DOCX Fixture Heading", level=1)
    document.add_paragraph("DOCX fixture paragraph for MarkItDown.")
    stream = BytesIO()
    document.save(stream)
    return stream.getvalue()


@pytest.fixture
def pdf_bytes() -> bytes:
    text = "PDF Fixture Heading"
    content = f"BT /F1 20 Tf 72 720 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content + b"\nendstream",
    ]

    pdf = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for index, body in enumerate(objects, start=1):
        offsets.append(len(pdf))
        pdf.extend(f"{index} 0 obj\n".encode())
        pdf.extend(body)
        pdf.extend(b"\nendobj\n")

    xref_offset = len(pdf)
    pdf.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    pdf.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf.extend(f"{offset:010d} 00000 n \n".encode())
    pdf.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode()
    )
    return bytes(pdf)


@pytest.fixture
def pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "PPTX Fixture Heading"
    slide.placeholders[1].text = "PPTX fixture paragraph for MarkItDown."
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


@pytest.fixture
def xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Fixture"
    sheet.append(["Name", "Value"])
    sheet.append(["XLSX Fixture", 42])
    stream = BytesIO()
    workbook.save(stream)
    return stream.getvalue()


@pytest.fixture
def html_bytes() -> bytes:
    return (
        b"<html><body><h1>HTML Fixture Heading</h1>"
        b"<p>Fixture paragraph.</p></body></html>"
    )
